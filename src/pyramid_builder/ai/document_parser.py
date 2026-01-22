"""
Document Parser for Strategic Content Extraction.

Parses PDF, DOCX, and PPTX files to extract text content.
"""

import io
from typing import Dict, List, Any
from pathlib import Path

try:
    from pypdf import PdfReader
    PDF_AVAILABLE = True
except ImportError:
    PDF_AVAILABLE = False

try:
    from docx import Document
    DOCX_AVAILABLE = True
except ImportError:
    DOCX_AVAILABLE = False

try:
    from pptx import Presentation
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False


class DocumentParser:
    """
    Parse various document formats to extract structured text content.
    Supports: PDF, DOCX, PPTX
    """

    MAX_FILE_SIZE = 10 * 1024 * 1024  # 10MB per file
    MAX_PAGES = 100

    @staticmethod
    def parse_pdf(file_content: bytes) -> Dict[str, Any]:
        """
        Parse PDF file and extract text content.

        Args:
            file_content: PDF file content as bytes

        Returns:
            Dict with extracted content
        """
        if not PDF_AVAILABLE:
            raise ImportError("pypdf not installed. Install with: pip install pypdf")

        if len(file_content) > DocumentParser.MAX_FILE_SIZE:
            raise ValueError(f"File size exceeds {DocumentParser.MAX_FILE_SIZE / 1024 / 1024}MB limit")

        try:
            pdf_file = io.BytesIO(file_content)
            reader = PdfReader(pdf_file)

            num_pages = len(reader.pages)
            if num_pages > DocumentParser.MAX_PAGES:
                raise ValueError(f"PDF has {num_pages} pages, exceeds {DocumentParser.MAX_PAGES} page limit")

            # Extract text from all pages
            text_blocks = []
            for page_num, page in enumerate(reader.pages, 1):
                text = page.extract_text()
                if text.strip():
                    text_blocks.append({
                        "page": page_num,
                        "content": text.strip(),
                        "type": "page"
                    })

            return {
                "success": True,
                "format": "pdf",
                "num_pages": num_pages,
                "blocks": text_blocks,
                "metadata": {
                    "title": reader.metadata.get("/Title", "") if reader.metadata else "",
                    "author": reader.metadata.get("/Author", "") if reader.metadata else "",
                }
            }

        except Exception as e:
            return {
                "success": False,
                "error": str(e),
                "format": "pdf"
            }

    @staticmethod
    def parse_docx(file_content: bytes) -> Dict[str, Any]:
        """
        Parse DOCX file and extract text content.

        Args:
            file_content: DOCX file content as bytes

        Returns:
            Dict with extracted content
        """
        if not DOCX_AVAILABLE:
            raise ImportError("python-docx not installed. Install with: pip install python-docx")

        if len(file_content) > DocumentParser.MAX_FILE_SIZE:
            raise ValueError(f"File size exceeds {DocumentParser.MAX_FILE_SIZE / 1024 / 1024}MB limit")

        try:
            docx_file = io.BytesIO(file_content)
            doc = Document(docx_file)

            # Extract paragraphs with structure
            text_blocks = []
            for para in doc.paragraphs:
                if para.text.strip():
                    # Detect if it's a heading based on style
                    is_heading = "Heading" in para.style.name if para.style else False

                    text_blocks.append({
                        "content": para.text.strip(),
                        "type": "heading" if is_heading else "paragraph",
                        "style": para.style.name if para.style else "Normal"
                    })

            # Also extract text from tables
            for table in doc.tables:
                for row in table.rows:
                    row_text = " | ".join([cell.text.strip() for cell in row.cells if cell.text.strip()])
                    if row_text:
                        text_blocks.append({
                            "content": row_text,
                            "type": "table_row"
                        })

            return {
                "success": True,
                "format": "docx",
                "blocks": text_blocks,
                "metadata": {
                    "num_paragraphs": len(doc.paragraphs),
                    "num_tables": len(doc.tables)
                }
            }

        except Exception as e:
            return {
                "success": False,
                "error": str(e),
                "format": "docx"
            }

    @staticmethod
    def parse_pptx(file_content: bytes) -> Dict[str, Any]:
        """
        Parse PPTX file and extract text content.

        Args:
            file_content: PPTX file content as bytes

        Returns:
            Dict with extracted content
        """
        if not PPTX_AVAILABLE:
            raise ImportError("python-pptx not installed. Install with: pip install python-pptx")

        if len(file_content) > DocumentParser.MAX_FILE_SIZE:
            raise ValueError(f"File size exceeds {DocumentParser.MAX_FILE_SIZE / 1024 / 1024}MB limit")

        try:
            pptx_file = io.BytesIO(file_content)
            prs = Presentation(pptx_file)

            num_slides = len(prs.slides)
            if num_slides > DocumentParser.MAX_PAGES:
                raise ValueError(f"Presentation has {num_slides} slides, exceeds {DocumentParser.MAX_PAGES} slide limit")

            # Extract text from slides
            text_blocks = []
            for slide_num, slide in enumerate(prs.slides, 1):
                slide_texts = []

                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        # Detect if it's a title
                        is_title = hasattr(shape, "is_placeholder") and shape.is_placeholder and shape.placeholder_format.type == 1

                        slide_texts.append({
                            "content": shape.text.strip(),
                            "type": "title" if is_title else "content"
                        })

                if slide_texts:
                    text_blocks.append({
                        "slide": slide_num,
                        "content": slide_texts,
                        "type": "slide"
                    })

            return {
                "success": True,
                "format": "pptx",
                "num_slides": num_slides,
                "blocks": text_blocks,
                "metadata": {
                    "num_slides": num_slides
                }
            }

        except Exception as e:
            return {
                "success": False,
                "error": str(e),
                "format": "pptx"
            }

    @classmethod
    def parse(cls, file_content: bytes, filename: str) -> Dict[str, Any]:
        """
        Parse document based on file extension.

        Args:
            file_content: File content as bytes
            filename: Original filename (used to determine format)

        Returns:
            Dict with extracted content
        """
        ext = Path(filename).suffix.lower()

        if ext == ".pdf":
            return cls.parse_pdf(file_content)
        elif ext == ".docx":
            return cls.parse_docx(file_content)
        elif ext == ".pptx":
            return cls.parse_pptx(file_content)
        else:
            return {
                "success": False,
                "error": f"Unsupported file format: {ext}. Supported formats: PDF, DOCX, PPTX",
                "format": ext
            }
