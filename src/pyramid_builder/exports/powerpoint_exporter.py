"""
PowerPoint (PPTX) export functionality for strategic pyramids.

Generates professional presentation slides with proper formatting.
"""

from typing import Optional
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor

from ..models.pyramid import StrategyPyramid


class PowerPointExporter:
    """Export pyramids to PowerPoint (PPTX) format with professional slides."""

    def __init__(self, pyramid: StrategyPyramid):
        """
        Initialize exporter.

        Args:
            pyramid: StrategyPyramid to export
        """
        self.pyramid = pyramid
        self.prs = Presentation()
        self.prs.slide_width = Inches(10)
        self.prs.slide_height = Inches(7.5)

        # Define color scheme
        self.primary_color = RGBColor(31, 119, 180)  # Blue
        self.secondary_color = RGBColor(100, 100, 100)  # Gray
        self.accent_color = RGBColor(255, 127, 14)  # Orange

    def export(
        self,
        filepath: str,
        audience: str = "leadership",
        include_title_slide: bool = True,
    ) -> Path:
        """
        Export pyramid to PowerPoint file.

        Args:
            filepath: Where to save the PowerPoint file
            audience: Target audience (executive, leadership, detailed)
            include_title_slide: Include a title slide

        Returns:
            Path to created file
        """
        if include_title_slide:
            self._add_title_slide()

        if audience == "executive":
            self._generate_executive_presentation()
        elif audience == "detailed":
            self._generate_detailed_presentation()
        else:  # leadership (default)
            self._generate_leadership_presentation()

        filepath_obj = Path(filepath)
        self.prs.save(str(filepath_obj))
        return filepath_obj

    def _add_title_slide(self):
        """Add professional title slide."""
        title_slide_layout = self.prs.slide_layouts[0]  # Title slide layout
        slide = self.prs.slides.add_slide(title_slide_layout)

        title = slide.shapes.title
        subtitle = slide.placeholders[1]

        title.text = self.pyramid.metadata.project_name
        subtitle.text = f"{self.pyramid.metadata.organization}\n{self.pyramid.metadata.created_by}"

        # Style the title
        title.text_frame.paragraphs[0].font.size = Pt(44)
        title.text_frame.paragraphs[0].font.bold = True
        title.text_frame.paragraphs[0].font.color.rgb = self.primary_color

    def _add_section_divider(self, title: str, subtitle: str = ""):
        """Add a section divider slide."""
        title_slide_layout = self.prs.slide_layouts[0]
        slide = self.prs.slides.add_slide(title_slide_layout)

        slide.shapes.title.text = title
        if subtitle:
            slide.placeholders[1].text = subtitle

        # Style
        slide.shapes.title.text_frame.paragraphs[0].font.size = Pt(54)
        slide.shapes.title.text_frame.paragraphs[0].font.color.rgb = self.primary_color

    def _add_content_slide(self, title: str, content_type: str = "bullet"):
        """Add a content slide with title."""
        if content_type == "bullet":
            layout = self.prs.slide_layouts[1]  # Title and Content
        else:
            layout = self.prs.slide_layouts[5]  # Title only
        slide = self.prs.slides.add_slide(layout)
        slide.shapes.title.text = title
        return slide

    def _generate_executive_presentation(self):
        """Generate executive presentation (5-8 slides)."""
        # Slide 1: Purpose
        if self.pyramid.vision:
            slide = self._add_content_slide("Our Purpose")
            text_frame = slide.placeholders[1].text_frame
            text_frame.clear()

            p = text_frame.paragraphs[0]
            p.text = self.pyramid.vision.statement
            p.font.size = Pt(24)
            p.font.italic = True
            p.alignment = PP_ALIGN.CENTER
            p.level = 0

        # Slide 2: Values
        if self.pyramid.values:
            slide = self._add_content_slide("Our Values")
            text_frame = slide.placeholders[1].text_frame
            text_frame.clear()

            for value in self.pyramid.values[:5]:  # Max 5 values
                p = text_frame.add_paragraph()
                p.text = value.name
                p.font.size = Pt(20)
                p.font.bold = True
                p.font.color.rgb = self.primary_color
                p.level = 0

                if value.description:
                    p2 = text_frame.add_paragraph()
                    p2.text = value.description
                    p2.font.size = Pt(16)
                    p2.level = 1

        # Slide 3: Strategic Drivers
        if self.pyramid.strategic_drivers:
            slide = self._add_content_slide("Strategic Focus")
            text_frame = slide.placeholders[1].text_frame
            text_frame.clear()

            for driver in self.pyramid.strategic_drivers:
                p = text_frame.add_paragraph()
                p.text = driver.name
                p.font.size = Pt(22)
                p.font.bold = True
                p.font.color.rgb = self.primary_color
                p.level = 0

                p2 = text_frame.add_paragraph()
                p2.text = driver.description
                p2.font.size = Pt(16)
                p2.level = 1

        # Slides 4-6: Key Commitments by Horizon
        if self.pyramid.iconic_commitments:
            for horizon in ["H1", "H2", "H3"]:
                commitments = [c for c in self.pyramid.iconic_commitments if c.horizon.value == horizon][:3]
                if commitments:
                    horizon_name = {
                        "H1": "Near-Term Commitments (H1)",
                        "H2": "Medium-Term Commitments (H2)",
                        "H3": "Long-Term Commitments (H3)"
                    }[horizon]

                    slide = self._add_content_slide(horizon_name)
                    text_frame = slide.placeholders[1].text_frame
                    text_frame.clear()

                    for commitment in commitments:
                        p = text_frame.add_paragraph()
                        p.text = commitment.name
                        p.font.size = Pt(18)
                        p.font.bold = True
                        p.level = 0

                        if commitment.target_date:
                            p2 = text_frame.add_paragraph()
                            p2.text = f"Target: {commitment.target_date}"
                            p2.font.size = Pt(14)
                            p2.font.italic = True
                            p2.level = 1

    def _generate_leadership_presentation(self):
        """Generate full leadership presentation (15-20 slides)."""
        # Section divider
        self._add_section_divider("Section 1: Purpose", "Why we exist")

        # Purpose slides
        if self.pyramid.vision:
            slide = self._add_content_slide("Our Vision")
            text_frame = slide.placeholders[1].text_frame
            text_frame.clear()

            p = text_frame.paragraphs[0]
            p.text = self.pyramid.vision.statement
            p.font.size = Pt(28)
            p.font.italic = True
            p.alignment = PP_ALIGN.CENTER

        if self.pyramid.values:
            slide = self._add_content_slide("Our Values")
            text_frame = slide.placeholders[1].text_frame
            text_frame.clear()

            for value in self.pyramid.values:
                p = text_frame.add_paragraph()
                p.text = value.name
                p.font.size = Pt(20)
                p.font.bold = True
                p.font.color.rgb = self.primary_color
                p.level = 0

                if value.description:
                    p2 = text_frame.add_paragraph()
                    p2.text = value.description
                    p2.font.size = Pt(16)
                    p2.level = 1

        # Section divider
        self._add_section_divider("Section 2: Strategy", "How we will succeed")

        # Behaviours
        if self.pyramid.behaviours:
            slide = self._add_content_slide("Our Behaviours")
            text_frame = slide.placeholders[1].text_frame
            text_frame.clear()

            for behaviour in self.pyramid.behaviours[:6]:  # Max 6 per slide
                p = text_frame.add_paragraph()
                p.text = behaviour.statement
                p.font.size = Pt(16)
                p.level = 0

        # Strategic Drivers (one slide per driver)
        for driver in self.pyramid.strategic_drivers:
            slide = self._add_content_slide(f"Strategic Driver: {driver.name}")
            text_frame = slide.placeholders[1].text_frame
            text_frame.clear()

            # Description
            p = text_frame.paragraphs[0]
            p.text = driver.description
            p.font.size = Pt(18)
            p.level = 0

            # Intents for this driver
            intents = [i for i in self.pyramid.strategic_intents if i.driver_id == driver.id]
            if intents:
                p_header = text_frame.add_paragraph()
                p_header.text = "\nWhat success looks like:"
                p_header.font.size = Pt(16)
                p_header.font.bold = True
                p_header.level = 0

                for intent in intents[:3]:  # Max 3 intents per slide
                    p_intent = text_frame.add_paragraph()
                    p_intent.text = intent.statement
                    p_intent.font.size = Pt(14)
                    p_intent.font.italic = True
                    p_intent.level = 1

        # Section divider
        self._add_section_divider("Section 3: Execution", "Our iconic commitments")

        # Commitments by horizon
        for horizon in ["H1", "H2", "H3"]:
            commitments = [c for c in self.pyramid.iconic_commitments if c.horizon.value == horizon]
            if commitments:
                horizon_name = {
                    "H1": "H1: Near-Term (0-12 months)",
                    "H2": "H2: Medium-Term (12-24 months)",
                    "H3": "H3: Long-Term (24-36 months)"
                }[horizon]

                slide = self._add_content_slide(horizon_name)
                text_frame = slide.placeholders[1].text_frame
                text_frame.clear()

                for commitment in commitments[:4]:  # Max 4 per slide
                    driver = self.pyramid.get_driver_by_id(commitment.primary_driver_id)
                    driver_name = driver.name if driver else "Unknown"

                    p = text_frame.add_paragraph()
                    p.text = commitment.name
                    p.font.size = Pt(18)
                    p.font.bold = True
                    p.level = 0

                    p2 = text_frame.add_paragraph()
                    details = f"Driver: {driver_name}"
                    if commitment.target_date:
                        details += f" | Target: {commitment.target_date}"
                    if commitment.owner:
                        details += f" | Owner: {commitment.owner}"
                    p2.text = details
                    p2.font.size = Pt(12)
                    p2.font.italic = True
                    p2.level = 1

        # Distribution slide
        if self.pyramid.iconic_commitments:
            slide = self._add_content_slide("Distribution Analysis")
            text_frame = slide.placeholders[1].text_frame
            text_frame.clear()

            distribution = self.pyramid.get_distribution_by_driver()
            total = sum(distribution.values())

            for driver_name, count in distribution.items():
                percentage = (count / total * 100) if total > 0 else 0
                p = text_frame.add_paragraph()
                p.text = f"{driver_name}: {count} commitments ({percentage:.0f}%)"
                p.font.size = Pt(18)
                p.level = 0

    def _generate_detailed_presentation(self):
        """Generate detailed presentation (25-30 slides)."""
        # Start with leadership presentation
        self._generate_leadership_presentation()

        # Add enablers slide
        if self.pyramid.enablers:
            slide = self._add_content_slide("Our Enablers")
            text_frame = slide.placeholders[1].text_frame
            text_frame.clear()

            for enabler in self.pyramid.enablers[:6]:
                p = text_frame.add_paragraph()
                enabler_type = f" ({enabler.enabler_type})" if enabler.enabler_type else ""
                p.text = f"{enabler.name}{enabler_type}"
                p.font.size = Pt(16)
                p.font.bold = True
                p.level = 0

                p2 = text_frame.add_paragraph()
                p2.text = enabler.description
                p2.font.size = Pt(12)
                p2.level = 1

        # Add team objectives if present
        if self.pyramid.team_objectives:
            self._add_section_divider("Team Objectives", "Departmental goals")

            # Group by team
            teams = {}
            for obj in self.pyramid.team_objectives:
                if obj.team_name not in teams:
                    teams[obj.team_name] = []
                teams[obj.team_name].append(obj)

            for team_name, objectives in teams.items():
                slide = self._add_content_slide(f"{team_name} Objectives")
                text_frame = slide.placeholders[1].text_frame
                text_frame.clear()

                for obj in objectives[:4]:  # Max 4 per slide
                    p = text_frame.add_paragraph()
                    p.text = obj.name
                    p.font.size = Pt(16)
                    p.font.bold = True
                    p.level = 0

                    if obj.metrics:
                        for metric in obj.metrics[:2]:  # Max 2 metrics
                            p2 = text_frame.add_paragraph()
                            p2.text = f"• {metric}"
                            p2.font.size = Pt(12)
                            p2.level = 1
